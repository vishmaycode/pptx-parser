import os
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ==================================================
# CONFIG
# ==================================================

PPTX_FILE = "input/sample.pptx"

# ==================================================
# HELPERS
# ==================================================

def emu_to_px(emu):
    return round(emu / 9525)

# ==================================================
# SETUP
# ==================================================

ppt_path = Path(PPTX_FILE)

ppt_name = ppt_path.stem

output_root = Path("output") / ppt_name
output_root.mkdir(parents=True, exist_ok=True)

prs = Presentation(PPTX_FILE)

# ==================================================
# PROCESS SLIDES
# ==================================================

for slide_index, slide in enumerate(prs.slides, start=1):

    print(f"Processing Slide {slide_index}")

    # ----------------------------------------------
    # Create slide folders
    # ----------------------------------------------

    slide_dir = output_root / f"slide_{slide_index}"
    images_dir = slide_dir / "images"

    slide_dir.mkdir(parents=True, exist_ok=True)
    images_dir.mkdir(parents=True, exist_ok=True)

    # ----------------------------------------------
    # Slide JSON structure
    # ----------------------------------------------

    slide_json = {
        "slide": slide_index,
        "width": emu_to_px(prs.slide_width),
        "height": emu_to_px(prs.slide_height),
        "elements": []
    }

    # ----------------------------------------------
    # Parse shapes
    # ----------------------------------------------

    for shape_index, shape in enumerate(slide.shapes, start=1):

        element = {
            "id": shape_index,
            "name": shape.name,
            "shape_type": str(shape.shape_type),

            # Position
            "x": emu_to_px(shape.left),
            "y": emu_to_px(shape.top),
            "width": emu_to_px(shape.width),
            "height": emu_to_px(shape.height)
        }

        # ------------------------------------------
        # TEXT
        # ------------------------------------------

        if hasattr(shape, "text") and shape.text.strip():

            element["type"] = "text"
            element["text"] = shape.text.strip()

            # Font extraction (first run only)
            try:
                paragraph = shape.text_frame.paragraphs[0]

                if paragraph.runs:
                    run = paragraph.runs[0]

                    if run.font:
                        element["font"] = {
                            "size": run.font.size.pt if run.font.size else None,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "name": run.font.name
                        }

            except Exception:
                pass

        # ------------------------------------------
        # IMAGE
        # ------------------------------------------

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

            image = shape.image

            ext = image.ext
            image_name = f"image_{shape_index}.{ext}"

            image_path = images_dir / image_name

            with open(image_path, "wb") as f:
                f.write(image.blob)

            element["type"] = "image"
            element["img_url"] = f"./images/{image_name}"

        # ------------------------------------------
        # TABLE
        # ------------------------------------------

        elif shape.has_table:

            table_data = []

            for row in shape.table.rows:

                row_data = []

                for cell in row.cells:
                    row_data.append(cell.text)

                table_data.append(row_data)

            element["type"] = "table"
            element["table"] = table_data

        # ------------------------------------------
        # OTHER SHAPES
        # ------------------------------------------

        else:
            element["type"] = "shape"

        slide_json["elements"].append(element)

    # ----------------------------------------------
    # Save meta.json
    # ----------------------------------------------

    meta_file = slide_dir / "meta.json"

    with open(meta_file, "w", encoding="utf-8") as f:
        json.dump(slide_json, f, indent=4, ensure_ascii=False)

# ==================================================
# DONE
# ==================================================

print("\nDone.")
print(f"Output saved to: {output_root}")